"""Add 'costs from actual bill' callout to the FinOps Engine deck.

Edits at the run level so existing fonts/colours are preserved.
Idempotent: running twice produces the same output.
"""
from pathlib import Path
from pptx import Presentation

DECK = Path(__file__).parent / "finops-engine-overview.pptx"


def get_run(slide, shape_idx, para_idx, run_idx):
    return slide.shapes[shape_idx].text_frame.paragraphs[para_idx].runs[run_idx]


def main() -> None:
    p = Presentation(DECK)
    slides = list(p.slides)

    # Slide 3 - footer line: append 'Costs from your actual bill' so the
    # message is anchored on the four-engines overview slide.
    get_run(slides[2], 20, 0, 0).text = (
        "Stdlib Python.   az login.   No agents.   No SaaS.   "
        "Read-only.   Costs from your actual bill, not list price."
    )

    # Slide 9 - caption beneath the example output: clarify that the
    # \u00a3/mo column comes from Cost Management actuals.
    get_run(slides[8], 15, 0, 0).text = (
        "Reply with `accept` / `defer` / `reject` per row.   "
        "\u00a3/mo = your Cost Management actuals (not list price)."
    )

    p.save(DECK)
    print(f"Updated {DECK}")


if __name__ == "__main__":
    main()
