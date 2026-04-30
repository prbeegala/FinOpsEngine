"""Refresh the FinOps Engine overview deck (v0.1.1 -> v0.1.3).

Edits at the *run* level so the deck's existing fonts, colours, and
sizing are preserved. Idempotent: running it twice produces the same
output.
"""
from pathlib import Path
from pptx import Presentation

DECK = Path(__file__).parent / "finops-engine-overview.pptx"
BULLET = "\u2022"
MIDDOT = "\u00b7"
ARROW = "\u2192"


def get_run(slide, shape_idx, para_idx, run_idx):
    return slide.shapes[shape_idx].text_frame.paragraphs[para_idx].runs[run_idx]


def main() -> None:
    p = Presentation(DECK)
    slides = list(p.slides)

    # Slide 1 - title footer version
    get_run(slides[0], 9, 0, 0).text = "v0.1.3"

    # Slide 4 - rightsizing-peak example
    get_run(slides[3], 11, 0, 0).text = (
        "python tools/rightsizing-peak/\n"
        "  rightsizing_peak.py `\n"
        "  --all-subs `\n"
        "  --days 30 `\n"
        "  --out-dir ./out/peak `\n"
        "  --downsize-cpu-p95-max 80 `\n"
        "  --upsize-cpu-p95-min   90"
    )

    # Slide 5 - hidden-waste example
    get_run(slides[4], 11, 0, 0).text = (
        "python tools/hidden-waste/\n"
        "  hidden_waste.py `\n"
        "  --all-subs `\n"
        "  --out-dir ./out/hidden\n\n"
        "# Outputs:\n"
        "#  hidden-waste-<date>.md\n"
        "#  hidden-waste-<date>.csv\n"
        "#  policy/*.audit.json"
    )

    # Slide 6 - ri-coverage example
    get_run(slides[5], 11, 0, 0).text = (
        "python tools/ri-coverage/\n"
        "  ri_coverage.py `\n"
        "  --all-subs `\n"
        "  --out-dir ./out/ri\n\n"
        "# Outputs:\n"
        "#  ri-coverage-<date>.md\n"
        "#  ri-shortlist-<date>.md"
    )

    # Slide 11 - Trust column: replace 'Unit tests + CI' (partially shipped in v0.1.3)
    get_run(slides[10], 15, 0, 0).text = f"{BULLET}  CI on PRs (pytest + py_compile)"

    # Slide 11 - version-train footer: add 'Shipped' line, keep 'Next' line
    get_run(slides[10], 20, 0, 0).text = (
        f"Shipped:  v0.1.2  --all-subs  {MIDDOT}  v0.1.3  fixture tests\n"
        f"Next:  v0.2.0  {ARROW}  low-risk wins   {MIDDOT}   "
        f"v0.3.0  {ARROW}  trust-the-automation bundle   {MIDDOT}   "
        f"v1.0.0  {ARROW}  schema lock + CI"
    )

    # Slide 12 - get-started example block
    get_run(slides[11], 5, 0, 0).text = (
        "git clone https://github.com/prbeegala/FinOpsEngine.git\n"
        "cd FinOpsEngine\n\n"
        "az login\n\n"
        "python tools/rightsizing-peak/rightsizing_peak.py `\n"
        "    --all-subs `\n"
        "    --days 30 `\n"
        "    --out-dir ./out/peak\n\n"
        "python tools/hidden-waste/hidden_waste.py `\n"
        "    --all-subs `\n"
        "    --out-dir ./out/hidden\n"
    )

    # Slide 12 - 'WHAT YOU NEED' last bullet: clarify --all-subs option
    get_run(slides[11], 8, 2, 0).text = (
        f"{BULLET}  A list of subscription IDs, or --all-subs for tenant-wide."
    )

    # Slide 12 - footer version
    get_run(slides[11], 11, 0, 0).text = (
        f"github.com/prbeegala/FinOpsEngine    {MIDDOT}    "
        f"Open source {MIDDOT} MIT {MIDDOT} v0.1.3"
    )

    p.save(DECK)
    print(f"Updated {DECK}")


if __name__ == "__main__":
    main()
