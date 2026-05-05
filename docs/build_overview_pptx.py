"""Build the FinOps Engine overview deck.

Run: python docs/build_overview_pptx.py
Outputs: docs/finops-engine-overview.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand palette — kept tight so the deck reads as one product.
NAVY    = RGBColor(0x0B, 0x1F, 0x3A)
TEAL    = RGBColor(0x00, 0x9B, 0x9E)
AMBER   = RGBColor(0xF2, 0xA1, 0x3B)
SLATE   = RGBColor(0x4A, 0x55, 0x68)
LIGHT   = RGBColor(0xF4, 0xF6, 0xF8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x10, 0x18, 0x28)


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back by inserting first.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, *,
                size=16, color=INK, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + b
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_band(slide, color, top=Inches(0), height=Inches(0.5)):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top,
                                  Inches(13.333), height)
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    return band


def add_pill(slide, left, top, width, height, label, *,
             fill=TEAL, color=WHITE, size=14, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return pill


def add_card(slide, left, top, width, height, title, body, *,
             accent=TEAL):
    """A title+body card with a coloured top stripe."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xE2, 0xE8)
    card.line.width = Pt(0.75)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, top, width, Inches(0.18))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    add_text(slide, left + Inches(0.2), top + Inches(0.28),
             width - Inches(0.3), Inches(0.5),
             title, size=18, bold=True, color=NAVY)
    add_bullets(slide, left + Inches(0.2), top + Inches(0.85),
                width - Inches(0.3), height - Inches(1),
                body, size=12, color=INK)


def add_footer(slide, left_text="FinOps Engine · Open source · MIT",
               right_text="github.com/prbeegala/FinOpsEngine"):
    add_text(slide, Inches(0.4), Inches(7.05),
             Inches(8), Inches(0.3),
             left_text, size=10, color=SLATE)
    add_text(slide, Inches(8.4), Inches(7.05),
             Inches(4.5), Inches(0.3),
             right_text, size=10, color=SLATE, align=PP_ALIGN.RIGHT)


def section_header(slide, eyebrow, title):
    add_band(slide, NAVY, top=Inches(0), height=Inches(0.5))
    add_band(slide, TEAL, top=Inches(0.5), height=Inches(0.06))
    add_text(slide, Inches(0.4), Inches(0.08),
             Inches(8), Inches(0.4),
             eyebrow, size=12, bold=True, color=WHITE)
    add_text(slide, Inches(0.4), Inches(0.75),
             Inches(12.5), Inches(0.6),
             title, size=28, bold=True, color=NAVY)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(s, NAVY)

    # Accent bars
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6),
                             Inches(2.0), Inches(0.18), Inches(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    add_text(s, Inches(1.0), Inches(1.7),
             Inches(11), Inches(0.5),
             "FINOPS ENGINE", size=14, bold=True, color=TEAL)

    add_text(s, Inches(1.0), Inches(2.2),
             Inches(11.5), Inches(2.0),
             "Peak-aware Azure cost optimisation,\nin four small engines.",
             size=44, bold=True, color=WHITE)

    add_text(s, Inches(1.0), Inches(4.4),
             Inches(11.5), Inches(0.6),
             "Beat Advisor's 7-day window.  Find the waste it can't price.  "
             "Own the buffer.",
             size=20, color=AMBER)

    add_pill(s, Inches(1.0), Inches(5.4), Inches(1.6), Inches(0.45),
             "Open source", fill=TEAL)
    add_pill(s, Inches(2.7), Inches(5.4), Inches(1.0), Inches(0.45),
             "MIT", fill=TEAL)
    add_pill(s, Inches(3.8), Inches(5.4), Inches(1.6), Inches(0.45),
             "Stdlib Python", fill=TEAL)
    add_pill(s, Inches(5.5), Inches(5.4), Inches(2.2), Inches(0.45),
             "az login · no agents", fill=TEAL)
    add_pill(s, Inches(7.8), Inches(5.4), Inches(0.9), Inches(0.45),
             "v0.1.1", fill=AMBER, color=NAVY)

    add_text(s, Inches(1.0), Inches(6.6),
             Inches(11.5), Inches(0.4),
             "github.com/prbeegala/FinOpsEngine",
             size=14, color=LIGHT)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, "THE PROBLEM",
                   "Today's tools leave money — and risk — on the table.")

    add_card(s, Inches(0.5), Inches(1.7), Inches(4.0), Inches(4.6),
             "Azure Advisor",
             ["7-day default window + 30-min max-of-avg buckets miss "
              "weekly / month-end / sub-30-min peaks.",
              "Ignores orphan waste (snapshots, empty plans, idle LBs).",
              "No owner routing. Findings die in a portal tab.",
              "No buffer-aware RI guidance."],
             accent=AMBER)

    add_card(s, Inches(4.7), Inches(1.7), Inches(4.0), Inches(4.6),
             "FinOps SaaS platforms",
             ["Charge a percentage of savings you find for them.",
              "Black-box scoring — hard to defend in a change review.",
              "Still hand-off remediation to your engineers.",
              "Vendor lock-in, data egress, security review tax."],
             accent=AMBER)

    add_card(s, Inches(8.9), Inches(1.7), Inches(4.0), Inches(4.6),
             "What's left to you",
             ["A weekly 90-minute spreadsheet walkthrough.",
              "Best-effort orphan hunts in Cost Management.",
              "RI buys based on a single-SKU portal screen.",
              "No memory of what was actioned vs deferred."],
             accent=AMBER)

    add_text(s, Inches(0.5), Inches(6.45),
             Inches(12.3), Inches(0.4),
             "Net: real money is saved by the practitioner, not the tool.",
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_answer(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, "THE ANSWER",
                   "Four small engines that do the work a senior practitioner does.")

    cards = [
        ("rightsizing-peak",
         ["30-day window + per-hour true-peak P95 / P99 "
          "(vs Advisor's 7-day / 30-min max-of-avg).",
          "Flags Advisor's unsafe downsizes.",
          "Tunable thresholds per run."],
         TEAL),
        ("hidden-waste",
         ["7 categories Advisor doesn't price.",
          "Real £ from Cost Mgmt actuals.",
          "Ships an Azure Policy starter pack."],
         AMBER),
        ("ri-coverage",
         ["Risk-scored RI / SP shortlist.",
          "Bounded by your refund buffer.",
          "Quantifies the buffer's binding cost."],
         TEAL),
        ("context-enricher",
         ["Joins findings to owner / criticality.",
          "One GitHub Issue per CODEOWNER.",
          "Replaces the weekly spreadsheet."],
         AMBER),
    ]
    x = Inches(0.5)
    for title, body, accent in cards:
        add_card(s, x, Inches(1.7), Inches(3.05), Inches(4.4),
                 title, body, accent=accent)
        x += Inches(3.20)

    add_text(s, Inches(0.5), Inches(6.3),
             Inches(12.3), Inches(0.4),
             "Stdlib Python.   az login.   No agents.   No SaaS.   "
             "Read-only by design.",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_engine(prs, *, eyebrow, title, lead, points, headline,
                 cli_example):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, eyebrow, title)

    add_text(s, Inches(0.5), Inches(1.65),
             Inches(7.6), Inches(0.6),
             lead, size=16, bold=True, color=SLATE)

    add_bullets(s, Inches(0.5), Inches(2.4),
                Inches(7.6), Inches(3.5),
                points, size=14, color=INK)

    # Headline callout
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(8.4), Inches(1.65),
                               Inches(4.5), Inches(2.0))
    panel.adjustments[0] = 0.05
    panel.fill.solid(); panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()
    add_text(s, Inches(8.6), Inches(1.85),
             Inches(4.2), Inches(0.4),
             "HEADLINE", size=11, bold=True, color=AMBER)
    add_text(s, Inches(8.6), Inches(2.25),
             Inches(4.2), Inches(1.4),
             headline, size=15, bold=True, color=WHITE)

    # CLI example
    code_panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(8.4), Inches(3.85),
                                    Inches(4.5), Inches(2.6))
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x32)
    code_panel.line.fill.background()
    add_text(s, Inches(8.55), Inches(3.95),
             Inches(4.2), Inches(0.35),
             "EXAMPLE", size=10, bold=True, color=AMBER, font="Consolas")
    add_text(s, Inches(8.55), Inches(4.25),
             Inches(4.25), Inches(2.1),
             cli_example, size=10, color=LIGHT, font="Consolas")
    add_footer(s)


def slide_rightsizing(prs):
    slide_engine(
        prs,
        eyebrow="ENGINE 1 / 4",
        title="rightsizing-peak — longer window, sharper peaks.",
        lead="Advisor already uses P95/P99, but on a 7-day default "
             "window with 30-min max-of-avg buckets. This engine uses "
             "30 days + per-hour true peaks.",
        points=[
            "30 days of per-hour Max CPU + Min memory from Azure Monitor.",
            "Per-VM verdict: DOWNSIZE / KEEP / UPSIZE / "
            "INSUFFICIENT_DATA.",
            "Cross-checks Advisor and flags any of its downsizes the "
            "engine deems unsafe.",
            "Seven CLI flags to tune downsize / upsize / coverage "
            "thresholds at runtime.",
            "Conservative / Balanced / Aggressive starter profiles in "
            "the README.",
        ],
        headline="Expect 1–5% of Advisor's downsize "
                 "recommendations to be unsafe.\nOne missed peak "
                 "wipes a year of savings.",
        cli_example=(
            "python tools/rightsizing-peak/\n"
            "  rightsizing_peak.py `\n"
            "  --subs <sub1>,<sub2> `\n"
            "  --days 30 `\n"
            "  --out-dir ./out/peak `\n"
            "  --downsize-cpu-p95-max 80 `\n"
            "  --upsize-cpu-p95-min   90"
        ),
    )


def slide_hidden_waste(prs):
    slide_engine(
        prs,
        eyebrow="ENGINE 2 / 4",
        title="hidden-waste — seven categories Advisor doesn't price.",
        lead="Finds the orphans, lifecycle gaps, and idle infrastructure "
             "Cost Management hides at the bottom of a list.",
        points=[
            "Empty App Service Plans · idle Standard load balancers.",
            "Orphan NICs · unattached premium disks (incl. ASR seed disks).",
            "Old snapshots · stopped-not-deallocated VMs · unused public IPs.",
            "Priced from 30-day Cost Management actuals; list-price "
            "fallback when actuals are missing.",
            "Ships a starter Azure Policy pack to prevent the same waste "
            "recurring.",
        ],
        headline="Single-digit-percent of monthly spend\nis the typical "
                 "first-pass result —\noften dominated by a couple of "
                 "outliers no human had spotted.",
        cli_example=(
            "python tools/hidden-waste/\n"
            "  hidden_waste.py `\n"
            "  --subs <sub1>,<sub2> `\n"
            "  --out-dir ./out/hidden\n\n"
            "# Outputs:\n"
            "#  hidden-waste-<date>.md\n"
            "#  hidden-waste-<date>.csv\n"
            "#  policy/*.audit.json"
        ),
    )


def slide_ri_coverage(prs):
    slide_engine(
        prs,
        eyebrow="ENGINE 3 / 4",
        title="ri-coverage — buy what fits the buffer, not the portal.",
        lead="Workload-aware Reservations & Compute Savings Plan shortlist, "
             "risk-scored against the cancellation-exposure buffer your "
             "procurement function actually has.",
        points=[
            "Per family × region demand curve, not single-SKU portal view.",
            "Defaults to a £5,000 refund buffer; overridable per run.",
            "Risk score per recommendation (HIGH / MED / LOW).",
            "Quantifies the additional savings unlocked by raising the "
            "buffer — usually the binding constraint.",
            "Cross-checked against rightsizing-peak so you never reserve "
            "a workload that should be downsized first.",
        ],
        headline="The binding constraint is usually\nprocurement, not "
                 "the data.\nPrice the buffer.",
        cli_example=(
            "python tools/ri-coverage/\n"
            "  ri_coverage.py `\n"
            "  --subs <sub1>,<sub2> `\n"
            "  --out-dir ./out/ri\n\n"
            "# Outputs:\n"
            "#  ri-coverage-<date>.md\n"
            "#  ri-shortlist-<date>.md"
        ),
    )


def slide_enricher(prs):
    slide_engine(
        prs,
        eyebrow="ENGINE 4 / 4",
        title="context-enricher — every finding has an owner by morning.",
        lead="Joins the engines' CSVs with criticality, owner, "
             "environment, and confidence — and routes findings to the "
             "right team automatically.",
        points=[
            "Owner resolved from CODEOWNERS (tag/YAML fallback on roadmap).",
            "Auto-issues HIGH and MED only; LOW is reported but not raised.",
            "One GitHub Issue per owner — edited in place each run, "
            "never duplicated.",
            "Per-row `accept` / `defer` / `reject` for closed-loop "
            "tracking.",
            "Replaces the recurring weekly FinOps spreadsheet "
            "walkthrough.",
        ],
        headline="One Issue per owner.\nOne morning.\nNo spreadsheet.",
        cli_example=(
            "python tools/context-enricher/\n"
            "  context_enricher.py `\n"
            "  --hidden-waste ./out/hidden `\n"
            "  --rightsizing  ./out/peak `\n"
            "  --out-dir ./out/enriched\n\n"
            "# Outputs:\n"
            "#  enriched-<date>.md\n"
            "#  issues/<owner>-<date>.md"
        ),
    )


def slide_how_it_runs(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, "HOW IT RUNS",
                   "One nightly workflow. ~10 minutes. Read-only.")

    # Pipeline boxes
    steps = [
        ("1. Auth",
         "GitHub OIDC →\nAzure Reader +\nCost Mgmt Reader",
         TEAL),
        ("2. Engines",
         "rightsizing-peak\nhidden-waste\nri-coverage",
         TEAL),
        ("3. Enrich",
         "context-enricher\njoins + routes\nto CODEOWNERS",
         AMBER),
        ("4. Publish",
         "Per-owner GitHub\nIssues · MD reports\n· CSV · Workbooks",
         AMBER),
    ]
    x = Inches(0.7)
    for i, (title, body, color) in enumerate(steps):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, Inches(2.2), Inches(2.5),
                                 Inches(2.4))
        box.adjustments[0] = 0.08
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)
        add_text(s, x + Inches(0.1), Inches(2.35),
                 Inches(2.3), Inches(0.5),
                 title, size=16, bold=True, color=NAVY)
        add_text(s, x + Inches(0.1), Inches(2.85),
                 Inches(2.3), Inches(1.5),
                 body, size=12, color=INK)
        if i < 3:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(2.55),
                                       Inches(3.15),
                                       Inches(0.4), Inches(0.5))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = SLATE
            arrow.line.fill.background()
        x += Inches(2.95)

    add_text(s, Inches(0.5), Inches(5.0),
             Inches(12.3), Inches(0.5),
             "What you get every morning",
             size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(5.5),
             Inches(12.3), Inches(1.4),
             "Markdown reports for execs   ·   "
             "CSVs for data teams   ·   "
             "Per-owner GitHub Issues for engineers   ·   "
             "Azure Monitor Workbooks for ops",
             size=14, color=INK, align=PP_ALIGN.CENTER)

    add_footer(s)


def slide_outputs(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, "WHAT THE OUTPUT LOOKS LIKE",
                   "Synthetic example — see /samples for the full set.")

    # Mock report panel
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(1.7),
                               Inches(7.5), Inches(5.0))
    panel.adjustments[0] = 0.02
    panel.fill.solid(); panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = RGBColor(0xCF, 0xD6, 0xDF)

    add_text(s, Inches(0.7), Inches(1.85),
             Inches(7.1), Inches(0.4),
             "## FinOps remediation queue — contoso-app-team",
             size=14, bold=True, color=NAVY, font="Consolas")
    add_text(s, Inches(0.7), Inches(2.25),
             Inches(7.1), Inches(0.4),
             "**12 findings · ~£8,420 / month (£101k / yr) recoverable.**",
             size=12, color=INK, font="Consolas")

    rows = [
        ("# | Resource          | Category    | Conf | £/mo"),
        ("--+-------------------+-------------+------+-----"),
        (" 1| asr-seed-disk-01  | unattached  | HIGH | 2,710"),
        (" 2| empty-asp-prod-04 | empty-asp   | HIGH | 1,290"),
        (" 3| snap-old-img-7    | old-snap    | HIGH |   840"),
        (" 4| pip-unused-eu-12  | unused-ip   | MED  |   190"),
        (" 5| nic-orphan-prod-3 | orphan-nic  | MED  |    35"),
        ("…  (truncated)"),
    ]
    y = Inches(2.75)
    for r in rows:
        add_text(s, Inches(0.7), y, Inches(7.1), Inches(0.32),
                 r, size=11, color=INK, font="Consolas")
        y += Inches(0.32)

    add_text(s, Inches(0.7), Inches(5.6),
             Inches(7.1), Inches(0.5),
             "Reply with `accept` / `defer` / `reject` per row.",
             size=11, color=SLATE, font="Consolas")

    # Right-hand: what consumers do with it
    add_text(s, Inches(8.4), Inches(1.85),
             Inches(4.5), Inches(0.4),
             "Who reads what",
             size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(8.4), Inches(2.4),
                Inches(4.5), Inches(4.5),
                ["Exec / FinOps lead → Markdown summary "
                 "(combined-*.md).",
                 "Data team → CSV (one row per resource, every column "
                 "you need to join).",
                 "Engineering owner → GitHub Issue, edited in place "
                 "each night.",
                 "Ops → Azure Monitor Workbook with the engine's "
                 "custom log table."],
                size=13)
    add_footer(s)


def slide_differentiators(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, "WHY IT'S DIFFERENT",
                   "Built for the practitioner, not the procurement deck.")

    rows = [
        ("",                       "Azure Advisor", "FinOps SaaS",
         "FinOps Engine"),
        ("Peak-aware verdicts",    "✗ Avg only",   "Sometimes",
         "✓ P95 / P99"),
        ("Prices orphan waste",    "✗",            "Partial",
         "✓ 7 categories"),
        ("Buffer-aware RI",        "✗",            "Some",
         "✓ Configurable"),
        ("Auto-routes to owners",  "✗",            "Premium tier",
         "✓ CODEOWNERS"),
        ("Open source · MIT",      "n/a",          "✗",
         "✓"),
        ("Cost model",             "Free",         "% of savings",
         "Free · stdlib"),
        ("Auditable verdicts",     "Black-box",    "Black-box",
         "Deterministic"),
    ]

    col_w = [Inches(3.6), Inches(3.0), Inches(3.0), Inches(3.0)]
    col_x = [Inches(0.5)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)
    y = Inches(1.85)
    row_h = Inches(0.55)

    for i, row in enumerate(rows):
        is_header = (i == 0)
        is_engine_col = False
        for j, cell in enumerate(row):
            box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     col_x[j], y, col_w[j], row_h)
            if is_header:
                box.fill.solid(); box.fill.fore_color.rgb = NAVY
                fg = WHITE
            elif j == 3:
                box.fill.solid()
                box.fill.fore_color.rgb = (RGBColor(0xE6, 0xF7, 0xF8)
                                           if i % 2 else
                                           RGBColor(0xD4, 0xEF, 0xF1))
                fg = NAVY
            else:
                box.fill.solid()
                box.fill.fore_color.rgb = (LIGHT if i % 2 else WHITE)
                fg = INK
            box.line.color.rgb = RGBColor(0xCF, 0xD6, 0xDF)
            box.line.width = Pt(0.5)

            tf = box.text_frame
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.05)
            p = tf.paragraphs[0]
            p.alignment = (PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
            run = p.add_run()
            run.text = cell
            run.font.size = Pt(13 if is_header else 12)
            run.font.bold = is_header or (j == 0 and not is_header) \
                            or (j == 3)
            run.font.color.rgb = fg
            run.font.name = "Calibri"
        y += row_h
    add_footer(s)


def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, "WHAT'S NEXT",
                   "Public roadmap — you can see the whole backlog on GitHub.")

    cards = [
        ("Coverage",
         ["Storage / Hot-tier waste",
          "Log Analytics & App Insights",
          "AKS node-pool peak rightsizing",
          "Dev / test auto-shutdown gap",
          "Cosmos DB autoscale",
          "Network · NAT / ER / VPN"],
         TEAL),
        ("Engine depth",
         ["Upsize + SKU-family swap",
          "State DB · trend tracking",
          "RI vs Savings Plan trade-off",
          "Tag / YAML owner routing"],
         TEAL),
        ("Trust the automation",
         ["Unit tests + CI",
          "Output schema versioning",
          "--plan-only dry-run",
          "Label-based dedupe",
          "Per-finding stable IDs"],
         AMBER),
        ("Reach",
         ["Docker · single CLI · YAML config",
          "Slack / Teams / ServiceNow / Jira",
          "Power BI dataset export",
          "AWS + GCP adapters",
          "M365 licensing waste"],
         AMBER),
    ]
    x = Inches(0.4)
    for title, body, accent in cards:
        add_card(s, x, Inches(1.7), Inches(3.1), Inches(4.7),
                 title, body, accent=accent)
        x += Inches(3.20)

    add_text(s, Inches(0.5), Inches(6.55),
             Inches(12.3), Inches(0.4),
             "v0.2.0  →  low-risk wins   ·   "
             "v0.3.0  →  trust-the-automation bundle   ·   "
             "v1.0.0  →  schema lock + tests + CI",
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_get_started(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, NAVY)
    add_band(s, TEAL, top=Inches(0), height=Inches(0.06))

    add_text(s, Inches(0.6), Inches(0.4),
             Inches(12), Inches(0.5),
             "GET STARTED IN 5 MINUTES", size=14, bold=True,
             color=TEAL)

    add_text(s, Inches(0.6), Inches(0.95),
             Inches(12.2), Inches(0.8),
             "Three commands. Real findings before lunch.",
             size=28, bold=True, color=WHITE)

    code = (
        "git clone https://github.com/prbeegala/FinOpsEngine.git\n"
        "cd FinOpsEngine\n\n"
        "az login\n"
        "az account set --subscription <default-sub-id>\n\n"
        "python tools/rightsizing-peak/rightsizing_peak.py `\n"
        "    --subs <sub1>,<sub2> `\n"
        "    --days 30 `\n"
        "    --out-dir ./out/peak-rightsizing\n\n"
        "python tools/hidden-waste/hidden_waste.py `\n"
        "    --subs <sub1>,<sub2> `\n"
        "    --out-dir ./out/hidden-waste\n"
    )

    code_panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.6), Inches(2.0),
                                    Inches(8.0), Inches(4.4))
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x32)
    code_panel.line.fill.background()
    add_text(s, Inches(0.8), Inches(2.15),
             Inches(7.6), Inches(4.1),
             code, size=13, color=LIGHT, font="Consolas")

    # CTA panel
    cta = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(8.9), Inches(2.0),
                             Inches(4.0), Inches(4.4))
    cta.adjustments[0] = 0.05
    cta.fill.solid(); cta.fill.fore_color.rgb = TEAL
    cta.line.fill.background()
    add_text(s, Inches(9.1), Inches(2.2),
             Inches(3.7), Inches(0.5),
             "WHAT YOU NEED", size=12, bold=True, color=NAVY)
    add_bullets(s, Inches(9.1), Inches(2.7),
                Inches(3.7), Inches(2.2),
                ["az login (Reader + Cost Mgmt Reader).",
                 "Python 3.10+ (no extra packages).",
                 "A list of subscription IDs."],
                size=13, color=WHITE)
    add_text(s, Inches(9.1), Inches(4.85),
             Inches(3.7), Inches(0.5),
             "DOCS", size=12, bold=True, color=NAVY)
    add_bullets(s, Inches(9.1), Inches(5.25),
                Inches(3.7), Inches(1.0),
                ["README.md · per-tool READMEs.",
                 "samples/ · synthetic outputs.",
                 "ROADMAP.md · what's coming."],
                size=12, color=WHITE)

    add_text(s, Inches(0.6), Inches(6.65),
             Inches(12.2), Inches(0.5),
             "github.com/prbeegala/FinOpsEngine    ·    Open source · MIT · v0.1.1",
             size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_answer(prs)
    slide_rightsizing(prs)
    slide_hidden_waste(prs)
    slide_ri_coverage(prs)
    slide_enricher(prs)
    slide_how_it_runs(prs)
    slide_outputs(prs)
    slide_differentiators(prs)
    slide_roadmap(prs)
    slide_get_started(prs)

    prs.save(str(out_path))
    print(f"Wrote {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    here = Path(__file__).resolve().parent
    build(here / "finops-engine-overview.pptx")
